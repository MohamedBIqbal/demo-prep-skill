#!/usr/bin/env python3
"""
Generate McKinsey-style PowerPoint Presentation

A template script for creating professional demo presentations.
Supports color palettes, screenshot integration, and 10-slide structure.

Usage:
    python generate_pptx.py --output my_demo.pptx
    python generate_pptx.py --output my_demo.pptx --screenshots ./screenshots
    python generate_pptx.py --output my_demo.pptx --palette iceberg

Requires:
    pip install python-pptx Pillow
"""

import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


# ============ COLOR PALETTES ============
# Multiple themes for different storytelling contexts

PALETTES = {
    'default': {
        'primary': RGBColor(0x00, 0x66, 0xCC),
        'primary_dark': RGBColor(0x00, 0x52, 0xA3),
        'secondary': RGBColor(0x1E, 0x29, 0x3B),
        'accent': RGBColor(0x10, 0xB9, 0x81),
        'warning': RGBColor(0xD9, 0x77, 0x06),
        'danger': RGBColor(0xDC, 0x26, 0x26),
        'success': RGBColor(0x16, 0xA3, 0x4A),
        'purple': RGBColor(0x8B, 0x5C, 0xF6),
        'white': RGBColor(0xFF, 0xFF, 0xFF),
        'light_gray': RGBColor(0xF8, 0xFA, 0xFC),
        'muted': RGBColor(0x64, 0x74, 0x8B),
        'border': RGBColor(0xE2, 0xE8, 0xF0),
    },
    # Iceberg extends default with storytelling colors for problem slides
    'iceberg': {
        # Base colors (same as default)
        'primary': RGBColor(0x00, 0x66, 0xCC),
        'primary_dark': RGBColor(0x00, 0x52, 0xA3),
        'secondary': RGBColor(0x1E, 0x29, 0x3B),
        'accent': RGBColor(0x10, 0xB9, 0x81),
        'warning': RGBColor(0xD9, 0x77, 0x06),
        'danger': RGBColor(0xDC, 0x26, 0x26),
        'success': RGBColor(0x16, 0xA3, 0x4A),
        'purple': RGBColor(0x8B, 0x5C, 0xF6),
        'white': RGBColor(0xFF, 0xFF, 0xFF),
        'light_gray': RGBColor(0xF8, 0xFA, 0xFC),
        'muted': RGBColor(0x64, 0x74, 0x8B),
        'border': RGBColor(0xE2, 0xE8, 0xF0),
        # Iceberg-specific storytelling colors
        'sky': RGBColor(0xE0, 0xF2, 0xFE),
        'sky_mid': RGBColor(0xBA, 0xE6, 0xFD),
        'ocean': RGBColor(0x02, 0x84, 0xC7),
        'ocean_deep': RGBColor(0x0C, 0x4A, 0x6E),
        'text_light': RGBColor(0xFF, 0xFF, 0xFF),
        'text_sky': RGBColor(0x03, 0x69, 0xA1),
    },
}

COLORS = PALETTES['default']

# Screenshot configuration - map slide numbers to expected file patterns
SCREENSHOT_SLIDES = {
    4: 'slide_04_architecture',
    5: 'slide_05_demo',
    6: 'slide_06_demo_edge',
    7: 'slide_07_proof',
    8: 'slide_08_audit',
}


# ============ SCREENSHOT HELPERS ============

def find_screenshot(screenshots_dir, slide_name, slide_num=None):
    """Find screenshot file for a slide (flexible naming, supports .png, .jpg, .jpeg)"""
    if not screenshots_dir:
        return None

    screenshots_path = Path(screenshots_dir)
    if not screenshots_path.exists():
        return None

    if slide_num:
        patterns = [
            f"*{slide_num:02d}*",
            f"*{slide_num}*",
        ]

        for pattern in patterns:
            for ext in ['.png', '.jpg', '.jpeg']:
                matches = list(screenshots_path.glob(f"{pattern}{ext}")) + \
                          list(screenshots_path.glob(f"{pattern}{ext.upper()}"))
                if matches:
                    return str(sorted(matches)[0])

    possible_names = [slide_name]
    if slide_num:
        possible_names.extend([
            f"slide_{slide_num:02d}",
            f"Slide_{slide_num:02d}",
            f"{slide_num:02d}",
            f"{slide_num}",
        ])

    for name in possible_names:
        for ext in ['.png', '.jpg', '.jpeg', '.PNG', '.JPG', '.JPEG']:
            screenshot_file = screenshots_path / f"{name}{ext}"
            if screenshot_file.exists():
                return str(screenshot_file)

    return None


def add_screenshot(slide, screenshot_path, x=0.75, y=2.3, max_width=11.5, max_height=4.2):
    """Add screenshot image to slide, preserving aspect ratio"""
    if screenshot_path and Path(screenshot_path).exists():
        from PIL import Image

        # Get original image dimensions
        with Image.open(screenshot_path) as img:
            img_width, img_height = img.size

        # Calculate aspect ratio
        aspect_ratio = img_width / img_height

        # Fit within max bounds while preserving aspect ratio
        if max_width / max_height > aspect_ratio:
            # Height is the constraint
            height = max_height
            width = height * aspect_ratio
        else:
            # Width is the constraint
            width = max_width
            height = width / aspect_ratio

        # Center horizontally within the max_width area
        x_centered = x + (max_width - width) / 2

        slide.shapes.add_picture(
            screenshot_path,
            Inches(x_centered), Inches(y), Inches(width), Inches(height)
        )
        return True
    return False


def add_screenshot_placeholder(slide, x=0.75, y=2.3, width=11.5, height=4.2, label="Add screenshot"):
    """Add placeholder for missing screenshot"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['light_gray']
    shape.line.color.rgb = COLORS['border']
    shape.line.width = Pt(2)
    shape.line.dash_style = 2  # Dashed

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[{label}]"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['muted']
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

    return shape


# ============ HELPER FUNCTIONS ============

def add_context_label(slide, text, color=None):
    """Add context label at top of slide (e.g., 'THE PROBLEM')"""
    ctx = slide.shapes.add_textbox(Inches(0.75), Inches(0.6), Inches(3), Inches(0.4))
    tf = ctx.text_frame
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = color or COLORS['primary']
    return ctx


def add_action_title(slide, text, y=1.0):
    """Add McKinsey action title - the takeaway, not a topic label"""
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(y), Inches(11.8), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['secondary']
    p.line_spacing = 1.2
    return title_box


def add_card(slide, x, y, width, height, title, content, icon="", title_color=None, border_color=None, bg_color=None):
    """Add a card with title and content"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color or COLORS['white']
    card.line.color.rgb = border_color or COLORS['border']
    card.line.width = Pt(1)

    title_y = y + 0.2
    title_box = slide.shapes.add_textbox(
        Inches(x + 0.2), Inches(title_y), Inches(width - 0.4), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon}  {title}" if icon else title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = title_color or COLORS['primary']

    content_box = slide.shapes.add_textbox(
        Inches(x + 0.2), Inches(title_y + 0.5), Inches(width - 0.4), Inches(height - 0.8)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['muted']
    p.line_spacing = 1.4

    return card


def add_speaker_notes(slide, text):
    """Add speaker notes to slide"""
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


# ============ SLIDE BUILDERS ============

def slide_1_title(prs, config):
    """Slide 1: Title - Your product and value prop"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(6), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = config.get('headline_part1', 'Your industry needs a')
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['secondary']

    p = tf.add_paragraph()
    p.text = config.get('headline_part2', 'better solution')
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    # Product name box
    brand_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.5), Inches(1.2), Inches(4.5), Inches(2.5)
    )
    brand_shape.fill.solid()
    brand_shape.fill.fore_color.rgb = COLORS['primary']
    brand_shape.line.fill.background()

    gx_box = slide.shapes.add_textbox(Inches(7.5), Inches(1.6), Inches(4.5), Inches(1.2))
    tf = gx_box.text_frame
    p = tf.paragraphs[0]
    p.text = config.get('product_name', 'Your Product')
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    subtitle = slide.shapes.add_textbox(Inches(7.5), Inches(2.9), Inches(4.5), Inches(0.5))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = config.get('tagline', 'YOUR TAGLINE HERE')
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
    p.alignment = PP_ALIGN.CENTER

    # Feature pillars
    pillars = config.get('pillars', [
        ("Feature 1", "Brief description"),
        ("Feature 2", "Brief description"),
        ("Feature 3", "Brief description"),
    ])
    pillar_x = 0.75
    for title, desc in pillars:
        add_card(slide, pillar_x, 4.5, 3.6, 1.8, title, desc)
        pillar_x += 3.9

    add_speaker_notes(slide, """TIMING: 30 seconds
SAY: Lead with your value prop
TRANSITION: Click next to show the problem""")

    return slide


def slide_2_problem_iceberg(prs, config):
    """Slide 2: The Problem - Iceberg Storytelling Design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    ICE = PALETTES['iceberg']

    add_context_label(slide, "The Problem", COLORS['danger'])
    add_action_title(slide, config.get('problem_title', 'State the core problem your audience faces'))

    # Left side: Requirements
    reqs = config.get('requirements', [
        ("Requirement 1", "Why this matters"),
        ("Requirement 2", "Why this matters"),
        ("Requirement 3", "Why this matters"),
    ])
    req_y = 2.6
    for title, desc in reqs:
        box = slide.shapes.add_textbox(Inches(0.75), Inches(req_y), Inches(5), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"  {title}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['secondary']
        p = tf.add_paragraph()
        p.text = f"     {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['muted']
        req_y += 0.8

    # Right side: Iceberg visualization
    iceberg_x = 6.3
    iceberg_width = 5.7

    # Sky (above water)
    sky = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(iceberg_x), Inches(2.4), Inches(iceberg_width), Inches(1.8)
    )
    sky.fill.solid()
    sky.fill.fore_color.rgb = ICE['sky']
    sky.line.fill.background()

    label1 = slide.shapes.add_textbox(Inches(iceberg_x), Inches(2.5), Inches(iceberg_width), Inches(0.3))
    tf = label1.text_frame
    p = tf.paragraphs[0]
    p.text = config.get('iceberg_above_label', 'THE VISIBLE COST').upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ICE['text_sky']
    p.alignment = PP_ALIGN.CENTER

    amount_box = slide.shapes.add_textbox(Inches(iceberg_x), Inches(2.85), Inches(iceberg_width), Inches(0.8))
    tf = amount_box.text_frame
    p = tf.paragraphs[0]
    p.text = config.get('iceberg_above_value', '$X.XM')
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xEA, 0x58, 0x0C)
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(iceberg_x), Inches(3.6), Inches(iceberg_width), Inches(0.3))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = config.get('iceberg_above_subtitle', "...but that's just the start")
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['muted']
    p.alignment = PP_ALIGN.CENTER

    # Ocean (below water)
    ocean = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(iceberg_x), Inches(4.2), Inches(iceberg_width), Inches(2.6)
    )
    ocean.fill.solid()
    ocean.fill.fore_color.rgb = ICE['ocean_deep']
    ocean.line.fill.background()

    label2 = slide.shapes.add_textbox(Inches(iceberg_x), Inches(4.35), Inches(iceberg_width), Inches(0.3))
    tf = label2.text_frame
    p = tf.paragraphs[0]
    p.text = config.get('iceberg_below_label', 'THE HIDDEN DAMAGE').upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ICE['sky_mid']
    p.alignment = PP_ALIGN.CENTER

    # Below water items
    below_items = config.get('iceberg_below_items', [
        {'title': 'Hidden cost 1', 'desc': 'Description'},
        {'title': 'Hidden cost 2', 'desc': 'Description'},
        {'title': 'Hidden cost 3', 'desc': 'Description'},
    ])
    item_y = 4.7
    for item in below_items[:3]:
        item_box = slide.shapes.add_textbox(Inches(iceberg_x + 0.3), Inches(item_y), Inches(iceberg_width - 0.6), Inches(0.55))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item['title']
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ICE['text_light']

        p = tf.add_paragraph()
        p.text = item['desc']
        p.font.size = Pt(10)
        p.font.color.rgb = ICE['sky_mid']

        item_y += 0.6

    add_speaker_notes(slide, """TIMING: 45 seconds
SAY: The visible cost is just the tip of the iceberg.
TRANSITION: "Here's the scale we're dealing with..." """)

    return slide


def slide_3_scale(prs, config):
    """Slide 3: Dataset/Problem Scale"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_context_label(slide, "The Scale")
    add_action_title(slide, config.get('scale_title', 'Show the magnitude of the problem'))

    subtitle = slide.shapes.add_textbox(Inches(0.75), Inches(1.9), Inches(11), Inches(0.5))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = config.get('scale_subtitle', 'Real data at production scale')
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['muted']

    stats = config.get('stats', [
        ('100+', 'metric 1', COLORS['primary'], True),
        ('50K', 'metric 2', COLORS['accent'], False),
        ('8', 'metric 3', COLORS['purple'], False),
    ])

    card_width = 3.7
    card_height = 2.2
    card_y = 2.8
    gap = 0.3
    card_x = 0.75

    for value, label, color, is_filled in stats:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(card_x), Inches(card_y), Inches(card_width), Inches(card_height)
        )
        card.fill.solid()
        if is_filled:
            card.fill.fore_color.rgb = color
            card.line.fill.background()
            text_color = COLORS['white']
            label_color = RGBColor(0xA0, 0xC4, 0xE8)
        else:
            card.fill.fore_color.rgb = COLORS['white']
            card.line.color.rgb = COLORS['border']
            card.line.width = Pt(2)
            text_color = color
            label_color = COLORS['muted']

        num_box = slide.shapes.add_textbox(Inches(card_x), Inches(card_y + 0.4), Inches(card_width), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER

        label_box = slide.shapes.add_textbox(Inches(card_x), Inches(card_y + 1.5), Inches(card_width), Inches(0.5))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = label_color
        p.alignment = PP_ALIGN.CENTER

        card_x += card_width + gap

    # Callout
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(5.3), Inches(11.5), Inches(0.9)
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(0xFF, 0xFB, 0xEB)
    callout.line.color.rgb = RGBColor(0xFC, 0xD3, 0x4D)
    callout.line.width = Pt(1)

    callout_text = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.6))
    tf = callout_text.text_frame
    p = tf.paragraphs[0]
    p.text = config.get('scale_callout', 'Example callout that emphasizes the scale')
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x92, 0x40, 0x0E)

    add_speaker_notes(slide, """TIMING: 30 seconds
SAY: Emphasize the magnitude
TRANSITION: "Here's the architecture..." """)

    return slide


def slide_4_architecture(prs, config, screenshots_dir=None):
    """Slide 4: Solution Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_context_label(slide, "Architecture")
    add_action_title(slide, config.get('architecture_title', 'How your solution works'))

    screenshot = find_screenshot(screenshots_dir, SCREENSHOT_SLIDES.get(4, ''), 4)
    if not add_screenshot(slide, screenshot):
        add_screenshot_placeholder(slide, label="Architecture diagram (4.png)")

    add_speaker_notes(slide, """TIMING: 45 seconds
SAY: Walk through the architecture
SHOW: Point to each stage""")
    return slide


def slide_5_demo_good(prs, config, screenshots_dir=None):
    """Slide 5: Demo - Happy Path"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_context_label(slide, "Live Demo", COLORS['accent'])
    add_action_title(slide, config.get('demo_good_title', 'Show the happy path'))

    screenshot = find_screenshot(screenshots_dir, SCREENSHOT_SLIDES.get(5, ''), 5)
    if not add_screenshot(slide, screenshot):
        add_screenshot_placeholder(slide, label="Demo screenshot - happy path (5.png)")

    add_speaker_notes(slide, """TIMING: 60 seconds
SAY: "Watch what happens when..."
SHOW: Run the demo""")
    return slide


def slide_6_demo_edge(prs, config, screenshots_dir=None):
    """Slide 6: Demo - Edge Case / Differentiator"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_context_label(slide, "Edge Case", COLORS['danger'])
    add_action_title(slide, config.get('demo_edge_title', 'Show how you handle edge cases'))

    screenshot = find_screenshot(screenshots_dir, SCREENSHOT_SLIDES.get(6, ''), 6)
    if not add_screenshot(slide, screenshot):
        add_screenshot_placeholder(slide, label="Demo screenshot - edge case (6.png)")

    add_speaker_notes(slide, """TIMING: 60 seconds
SAY: "This is the key differentiator"
SHOW: Edge case handling""")
    return slide


def slide_7_proof(prs, config, screenshots_dir=None):
    """Slide 7: Results / Proof"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_context_label(slide, "The Proof", COLORS['accent'])
    add_action_title(slide, config.get('proof_title', 'Metrics and evidence'))

    screenshot = find_screenshot(screenshots_dir, SCREENSHOT_SLIDES.get(7, ''), 7)
    if not add_screenshot(slide, screenshot, max_height=3.5):
        # Show metric cards instead
        metrics = config.get('metrics', [
            ('95%', 'Metric 1'),
            ('2.5x', 'Metric 2'),
            ('100%', 'Metric 3'),
        ])
        card_x = 0.75
        for value, label in metrics:
            add_card(slide, card_x, 2.8, 3.6, 2.0, f"✓ {label}", "", title_color=COLORS['accent'])
            num_box = slide.shapes.add_textbox(Inches(card_x + 0.2), Inches(3.4), Inches(3.2), Inches(0.8))
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = COLORS['secondary']
            card_x += 3.9

    add_speaker_notes(slide, """TIMING: 30 seconds
SAY: Back up your claims with numbers""")
    return slide


def slide_8_audit(prs, config, screenshots_dir=None):
    """Slide 8: Audit Trail / Traceability (optional)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_context_label(slide, "Traceability", COLORS['purple'])
    add_action_title(slide, config.get('audit_title', 'Full audit trail and logging'))

    screenshot = find_screenshot(screenshots_dir, SCREENSHOT_SLIDES.get(8, ''), 8)
    if not add_screenshot(slide, screenshot):
        add_screenshot_placeholder(slide, label="Audit trail view (8.png)")

    add_speaker_notes(slide, """TIMING: 30 seconds
SAY: Complete traceability""")
    return slide


def slide_9_roadmap(prs, config):
    """Slide 9: Roadmap + Honest Gaps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_context_label(slide, "Roadmap")
    add_action_title(slide, config.get('roadmap_title', "What's done and what's next"))

    # Left: Progress
    progress_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(2.5), Inches(5.3), Inches(4)
    )
    progress_shape.fill.solid()
    progress_shape.fill.fore_color.rgb = COLORS['white']
    progress_shape.line.color.rgb = COLORS['border']

    prog_header = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(4.8), Inches(0.5))
    tf = prog_header.text_frame
    p = tf.paragraphs[0]
    p.text = "Progress"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    items = config.get('roadmap_items', [
        ("✓", "Completed item 1", "Details", COLORS['success']),
        ("2", "Next milestone", "Details", COLORS['primary']),
        ("3", "Future goal", "Details", COLORS['primary']),
    ])
    item_y = 3.3
    for num, title, desc, color in items:
        box = slide.shapes.add_textbox(Inches(1), Inches(item_y), Inches(4.8), Inches(0.7))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{num}   {title}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p = tf.add_paragraph()
        p.text = f"     {desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['muted']
        item_y += 0.8

    # Right: Gaps
    gaps_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.4), Inches(2.5), Inches(5.3), Inches(4)
    )
    gaps_shape.fill.solid()
    gaps_shape.fill.fore_color.rgb = COLORS['white']
    gaps_shape.line.color.rgb = COLORS['border']

    gaps_header = slide.shapes.add_textbox(Inches(6.65), Inches(2.7), Inches(4.8), Inches(0.5))
    tf = gaps_header.text_frame
    p = tf.paragraphs[0]
    p.text = "Honest Gaps"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['warning']

    gaps = config.get('gaps', [
        ("1", "Known limitation", "Context"),
        ("2", "Future work", "Context"),
        ("3", "Open question", "Context"),
    ])
    gap_y = 3.3
    for num, title, desc in gaps:
        box = slide.shapes.add_textbox(Inches(6.65), Inches(gap_y), Inches(4.8), Inches(0.7))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{num}   {title}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLORS['warning']
        p = tf.add_paragraph()
        p.text = f"     {desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['muted']
        gap_y += 0.8

    add_speaker_notes(slide, """TIMING: 30 seconds
SAY: Be honest about gaps
SHOW: Point to roadmap""")

    return slide


def slide_10_ask(prs, config):
    """Slide 10: The Ask"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_context_label(slide, "The Ask")
    add_action_title(slide, config.get('ask_title', 'What I need from you'))

    # Question cards
    q1_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(2.3), Inches(5.3), Inches(1.5)
    )
    q1_shape.fill.solid()
    q1_shape.fill.fore_color.rgb = COLORS['white']
    q1_shape.line.color.rgb = COLORS['border']

    q1 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(4.8), Inches(1.2))
    tf = q1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = config.get('ask_1_title', 'Feedback Request')
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p = tf.add_paragraph()
    p.text = config.get('ask_1_desc', 'What specific feedback do you want?')
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['muted']

    q2_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.4), Inches(2.3), Inches(5.3), Inches(1.5)
    )
    q2_shape.fill.solid()
    q2_shape.fill.fore_color.rgb = COLORS['white']
    q2_shape.line.color.rgb = COLORS['border']

    q2 = slide.shapes.add_textbox(Inches(6.65), Inches(2.5), Inches(4.8), Inches(1.2))
    tf = q2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = config.get('ask_2_title', 'Priority Question')
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['warning']
    p = tf.add_paragraph()
    p.text = config.get('ask_2_desc', 'What decision do you need help with?')
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['muted']

    # Thank you
    thanks = slide.shapes.add_textbox(Inches(0.75), Inches(5.5), Inches(11.5), Inches(0.8))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    questions = slide.shapes.add_textbox(Inches(0.75), Inches(6.2), Inches(11.5), Inches(0.4))
    tf = questions.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['muted']
    p.alignment = PP_ALIGN.CENTER

    add_speaker_notes(slide, """TIMING: 30 seconds
SAY: Be specific about what you need
ASK: Open it up for questions""")

    return slide


# ============ MAIN ============

def generate_presentation(output_path, config=None, screenshots_dir=None, palette='default'):
    """Generate the full 10-slide presentation"""
    global COLORS
    COLORS = PALETTES.get(palette, PALETTES['default'])

    if config is None:
        config = {}

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # Check screenshots
    screenshots_found = []
    screenshots_missing = []
    if screenshots_dir and Path(screenshots_dir).exists():
        print(f"✓ Screenshots directory: {screenshots_dir}")
        for slide_num, slide_name in SCREENSHOT_SLIDES.items():
            if find_screenshot(screenshots_dir, slide_name, slide_num):
                screenshots_found.append(slide_num)
            else:
                screenshots_missing.append(slide_num)

    # Build all 10 slides
    slide_1_title(prs, config)
    slide_2_problem_iceberg(prs, config)
    slide_3_scale(prs, config)
    slide_4_architecture(prs, config, screenshots_dir)
    slide_5_demo_good(prs, config, screenshots_dir)
    slide_6_demo_edge(prs, config, screenshots_dir)
    slide_7_proof(prs, config, screenshots_dir)
    slide_8_audit(prs, config, screenshots_dir)
    slide_9_roadmap(prs, config)
    slide_10_ask(prs, config)

    prs.save(output_path)
    print(f"✓ Generated: {output_path}")
    print(f"  10 slides with {palette} palette")

    if screenshots_found:
        print(f"  ✓ Screenshots inserted for slides: {screenshots_found}")
    if screenshots_missing:
        print(f"  ⚠ Screenshots missing for slides: {screenshots_missing}")
        print(f"    Add to: {screenshots_dir}/")


def main():
    parser = argparse.ArgumentParser(description="Generate Demo PPTX")
    parser.add_argument("--output", "-o", default="demo.pptx", help="Output file path")
    parser.add_argument("--screenshots", "-s", default=None, help="Path to screenshots directory")
    parser.add_argument("--palette", "-p", default="default", choices=['default', 'iceberg'],
                        help="Color palette to use")
    args = parser.parse_args()

    # ============ CUSTOMIZE YOUR CONTENT HERE ============
    config = {
        'product_name': 'Your Product',
        'tagline': 'YOUR TAGLINE HERE',
        'headline_part1': 'Your industry needs a',
        'headline_part2': 'better solution',
        'pillars': [
            ('Feature 1', 'Brief description'),
            ('Feature 2', 'Brief description'),
            ('Feature 3', 'Brief description'),
        ],
        'problem_title': 'State the core problem — be specific',
        'requirements': [
            ('Requirement 1', 'Why this matters'),
            ('Requirement 2', 'Why this matters'),
            ('Requirement 3', 'Why this matters'),
        ],
        'iceberg_above_label': 'The visible cost',
        'iceberg_above_value': '$X.XM',
        'iceberg_above_subtitle': "...but that's just the start",
        'iceberg_below_label': 'The hidden damage',
        'iceberg_below_items': [
            {'title': 'Hidden cost 1', 'desc': 'No audit trail = no answer'},
            {'title': 'Hidden cost 2', 'desc': 'Triggered by missing records'},
            {'title': 'Hidden cost 3', 'desc': 'Wrong answer → wrong decision'},
        ],
        'scale_title': 'Built on real data at production scale',
        'scale_subtitle': 'Actual data — indexed and ready',
        'stats': [
            ('100+', 'documents', PALETTES['default']['primary'], True),
            ('50K', 'searchable items', PALETTES['default']['accent'], False),
            ('8', 'categories', PALETTES['default']['purple'], False),
        ],
        'scale_callout': 'Example: A single document can be 1,000+ pages',
        'architecture_title': 'How it works — architecture overview',
        'demo_good_title': 'A well-grounded result',
        'demo_edge_title': 'Handling edge cases — the key differentiator',
        'proof_title': 'Metrics that prove it works',
        'metrics': [('95%', 'Metric 1'), ('2.5x', 'Metric 2'), ('100%', 'Metric 3')],
        'audit_title': 'Complete decision lineage — every step logged',
        'roadmap_title': "What's done — and what's next",
        'roadmap_items': [
            ("✓", "Core functionality", "Validated with tests", PALETTES['default']['success']),
            ("2", "Next milestone", "In progress", PALETTES['default']['primary']),
            ("3", "Future goal", "Planned", PALETTES['default']['primary']),
        ],
        'gaps': [
            ("1", "Known limitation", "Context here"),
            ("2", "Future work", "Context here"),
            ("3", "Open question", "Context here"),
        ],
        'ask_title': 'Your feedback — and guidance on priorities',
        'ask_1_title': 'Feedback on Implementation',
        'ask_1_desc': 'Is this approach sound?',
        'ask_2_title': 'What to Prioritize Next?',
        'ask_2_desc': 'Which feature should come first?',
    }

    generate_presentation(args.output, config, args.screenshots, args.palette)


if __name__ == "__main__":
    main()
